#!/usr/bin/env python3
"""Assemble presentation deck from processed metrics and figure PNGs (Part 5)."""

from __future__ import annotations

import os
import sys
from dataclasses import dataclass
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

_ROOT = Path(__file__).resolve().parents[1]
_PROCESSED = _ROOT / "data" / "processed"
_FIGURES = _ROOT / "outputs" / "figures"
_OUTPUT = _ROOT / "outputs" / "Polymarket_Growth_Teardown.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette (matches chart theme)
COLOR_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_POLITICS = RGBColor(0xD1, 0x49, 0x5B)
COLOR_SPORTS = RGBColor(0x3A, 0x6E, 0xA5)
COLOR_CRYPTO = RGBColor(0xE8, 0xA3, 0x3D)
COLOR_OTHER = RGBColor(0x9A, 0xA0, 0xA6)
COLOR_BOX_BG = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_BOX_BORDER = RGBColor(0xDD, 0xDD, 0xDD)

FONT = "Calibri"
DECK_AUTHOR = os.environ.get("DECK_AUTHOR", "Your Name")

SPIKE_MONTHS = frozenset({"2024-10", "2025-10", "2026-03"})


@dataclass(frozen=True)
class DeckMetrics:
    """All headline numbers loaded from data/processed/*.csv."""

    total_users: int
    politics_users: int
    politics_pct: float
    sports_users: int
    crypto_users: int
    other_users: int
    baseline_lo: int
    baseline_hi: int
    spike_election: int
    spike_election_month: str
    spike_recent: int
    spike_recent_month: str
    ever_returned_pooled: float
    ever_returned_median: float
    claimed_return: float
    seq_m1: float
    seq_m3: float
    seq_m6: float
    churn_m6: float
    m12_politics: float
    m12_sports: float
    m12_crypto: float
    m12_other: float
    m12_lo: float
    m12_hi: float
    election_m12_politics: float
    election_m12_sports: float
    election_m12_crypto: float
    north_star_m6: float
    lift_5pp_on_200k: int


def _pct(value: float, decimals: int = 1) -> str:
    return f"{value * 100:.{decimals}f}%"


def _k(value: int) -> str:
    if value >= 1_000_000:
        return f"{value / 1_000_000:.2f}M"
    return f"{value / 1_000:.0f}k"


def load_metrics() -> DeckMetrics:
    mix = pd.read_csv(_PROCESSED / "acquisition_mix.csv")
    acq = pd.read_csv(_PROCESSED / "acquisition_by_month.csv")
    ever = pd.read_csv(_PROCESSED / "ever_returned_summary.csv").iloc[0]
    seq = pd.read_csv(_PROCESSED / "median_sequential_retention.csv", index_col=0)
    mature = pd.read_csv(_PROCESSED / "mature_average_retention.csv", index_col=0)
    election = pd.read_csv(_PROCESSED / "election_cohort_retention.csv", index_col=0)

    total_users = int(mix["acquired_users"].sum())
    politics_row = mix[mix["first_category"] == "politics"].iloc[0]

    steady = acq[
        ~acq["cohort_label"].isin(SPIKE_MONTHS)
        & (acq["cohort_size"] >= 30_000)
        & (acq["cohort_size"] <= 50_000)
    ]
    if steady.empty:
        steady = acq[~acq["cohort_label"].isin(SPIKE_MONTHS)]

    election_spike = acq[acq["cohort_label"] == "2024-10"].iloc[0]
    recent_spike = acq[acq["cohort_label"] == "2026-03"].iloc[0]

    m12 = mature.loc[12]
    seq_m6 = float(seq.loc[6, "median_retention_rate"])

    return DeckMetrics(
        total_users=total_users,
        politics_users=int(politics_row["acquired_users"]),
        politics_pct=float(politics_row["pct_of_users"]) / 100.0,
        sports_users=int(mix[mix["first_category"] == "sports"]["acquired_users"].iloc[0]),
        crypto_users=int(mix[mix["first_category"] == "crypto"]["acquired_users"].iloc[0]),
        other_users=int(mix[mix["first_category"] == "other"]["acquired_users"].iloc[0]),
        baseline_lo=int(steady["cohort_size"].min()),
        baseline_hi=int(steady["cohort_size"].max()),
        spike_election=int(election_spike["cohort_size"]),
        spike_election_month="2024-10",
        spike_recent=int(recent_spike["cohort_size"]),
        spike_recent_month="2026-03",
        ever_returned_pooled=float(ever["pooled_rate"]),
        ever_returned_median=float(ever["median_rate"]),
        claimed_return=0.75,
        seq_m1=float(seq.loc[1, "median_retention_rate"]),
        seq_m3=float(seq.loc[3, "median_retention_rate"]),
        seq_m6=seq_m6,
        churn_m6=1.0 - seq_m6,
        m12_politics=float(m12["politics"]),
        m12_sports=float(m12["sports"]),
        m12_crypto=float(m12["crypto"]),
        m12_other=float(m12["other"]),
        m12_lo=float(m12.min()),
        m12_hi=float(m12.max()),
        election_m12_politics=float(election.loc[12, "politics"]),
        election_m12_sports=float(election.loc[12, "sports"]),
        election_m12_crypto=float(election.loc[12, "crypto"]),
        north_star_m6=seq_m6,
        lift_5pp_on_200k=int(200_000 * 0.05),
    )


def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def _textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.2,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = line_spacing
    return box


def _add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    *,
    font_size: int = 20,
    spacing: float = 1.35,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT
        p.level = 0
        p.line_spacing = spacing
        p.space_after = Pt(10)
    return box


def _speaker_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def _add_image(slide, path: Path, left, top, width, height) -> None:
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def _headline(slide, text: str) -> None:
    _textbox(
        slide,
        Inches(0.75),
        Inches(0.55),
        Inches(11.8),
        Inches(1.1),
        text,
        font_size=28,
        bold=True,
    )


def _subtext(slide, text: str) -> None:
    _textbox(
        slide,
        Inches(0.75),
        Inches(1.65),
        Inches(11.8),
        Inches(0.7),
        text,
        font_size=14,
        color=COLOR_MUTED,
    )


def _footer(slide, text: str) -> None:
    _textbox(
        slide,
        Inches(0.75),
        Inches(7.05),
        Inches(11.8),
        Inches(0.35),
        text,
        font_size=10,
        color=COLOR_MUTED,
    )


def slide_title(prs: Presentation, m: DeckMetrics) -> None:
    slide = _blank_slide(prs)
    _textbox(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(11.5),
        Inches(1.4),
        "Polymarket Growth: A Cohort Retention Teardown",
        font_size=36,
        bold=True,
        align=PP_ALIGN.LEFT,
    )
    _textbox(
        slide,
        Inches(0.9),
        Inches(3.65),
        Inches(11.0),
        Inches(0.9),
        f"What {_k(m.total_users)} traders' on-chain history says about where growth actually leaks",
        font_size=20,
        color=COLOR_MUTED,
    )
    _footer(slide, f"On-chain Polygon data via Dune · {DECK_AUTHOR}")
    _speaker_notes(
        slide,
        "Open with the scope: this is a demand-side teardown using on-chain taker activity, "
        f"not a product funnel audit. We tracked {_k(m.total_users)} unique taker wallets from "
        "Jan 2024 onward. The question for Growth: where does value leak after acquisition spikes?",
    )


def slide_bluf(prs: Presentation, m: DeckMetrics) -> None:
    slide = _blank_slide(prs)
    _headline(slide, "Bottom line up front")
    bullets = [
        "Acquisition is solved — and violently event-driven.",
        (
            f"The leak is retention: {_pct(m.churn_m6, 0)} of every cohort is gone within "
            f"6 months (median sequential M6 = {_pct(m.seq_m6)})."
        ),
        (
            f"Most users return once ({_pct(m.ever_returned_pooled)} ever-returned, pooled mature "
            f"2024–2025; near Polymarket's claimed {_pct(m.claimed_return, 0)}) — but never form a habit "
            f"({_pct(m.seq_m1)} → {_pct(m.seq_m3)} → {_pct(m.seq_m6)} sequential by M6)."
        ),
        (
            f"Recommendation: build the retention/habit layer. North Star = month-6 sequential "
            f"retention ({_pct(m.north_star_m6)} today)."
        ),
    ]
    _add_bullets(slide, Inches(0.9), Inches(1.85), Inches(11.5), Inches(4.5), bullets, font_size=22)
    _speaker_notes(
        slide,
        "Four beats: (1) Polymarket wins acquisition during mega-events — don't fight that. "
        f"(2) The durable problem is retention — {_pct(m.churn_m6, 0)} churn by M6 on the sequential "
        f"metric. (3) The '~75% return' claim is partly defensible on an ever-returned basis "
        f"({_pct(m.ever_returned_pooled)} in our data) but misleading if read as habit — sequential "
        "retention collapses. (4) Invest in habit formation, not just top-of-funnel.",
    )


def slide_acquisition(prs: Presentation, m: DeckMetrics) -> None:
    slide = _blank_slide(prs)
    _headline(slide, "Acquisition is violently event-driven — and Polymarket is great at it")
    _subtext(
        slide,
        f"~{_k(m.baseline_lo)}–{_k(m.baseline_hi)} new traders/mo baseline, spiking to "
        f"{_k(m.spike_election)} ({m.spike_election_month} election) and "
        f"{_k(m.spike_recent)} ({m.spike_recent_month}).",
    )
    _add_image(
        slide,
        _FIGURES / "01_acquisition_by_month.png",
        Inches(0.75),
        Inches(2.35),
        Inches(11.85),
        Inches(4.55),
    )
    _speaker_notes(
        slide,
        f"Chart shows monthly new taker cohorts from {m.baseline_lo:,}–{m.baseline_hi:,} in steady "
        f"months to {m.spike_election:,} in Oct 2024 and {m.spike_recent:,} in Mar 2026. "
        "Growth's job during spikes is not to dampen volume — it's to convert event attention "
        "into durable traders. The baseline proves organic demand exists between events.",
    )


def slide_75pct(prs: Presentation, m: DeckMetrics) -> None:
    slide = _blank_slide(prs)
    _headline(slide, "Most users come back once — they just don't stick")
    _subtext(
        slide,
        f"{_pct(m.ever_returned_pooled)} ever return (pooled mature vs claimed {_pct(m.claimed_return, 0)}), "
        f"but sequential retention collapses {_pct(m.seq_m1)} → {_pct(m.seq_m3)} → {_pct(m.seq_m6)} "
        f"by month 6. The habit gap is the opportunity.",
    )
    _add_image(
        slide,
        _FIGURES / "04_the_75pct_myth.png",
        Inches(0.75),
        Inches(2.35),
        Inches(11.85),
        Inches(4.55),
    )
    _speaker_notes(
        slide,
        "This slide hardens the '~75% return' claim. Ever-returned (≥2 active calendar months) "
        f"pools to {_pct(m.ever_returned_pooled)} across mature 2024–2025 cohorts — closer to "
        f"Polymarket's {_pct(m.claimed_return, 0)} than our sequential metric. But sequential "
        f"month-N retention (active in exactly month N) is stricter: median {_pct(m.seq_m1)} at M1, "
        f"{_pct(m.seq_m3)} at M3, {_pct(m.seq_m6)} at M6. Users visit twice; they don't build "
        "a monthly habit. That's the product gap.",
    )


def slide_category(prs: Presentation, m: DeckMetrics) -> None:
    slide = _blank_slide(prs)
    _headline(slide, "You can't fix it by acquiring 'better' users — entry category barely matters")
    _subtext(
        slide,
        f"All four categories converge to a {_pct(m.m12_lo)}–{_pct(m.m12_hi)} band by month 12 "
        f"(mature cohorts: politics {_pct(m.m12_politics)}, sports {_pct(m.m12_sports)}, "
        f"crypto {_pct(m.m12_crypto)}, other {_pct(m.m12_other)}).",
    )
    _add_image(
        slide,
        _FIGURES / "02_retention_by_category.png",
        Inches(0.75),
        Inches(2.35),
        Inches(11.85),
        Inches(4.55),
    )
    _speaker_notes(
        slide,
        "Mature average retention curves (2024-05..2025-05 cohorts, mean-of-rates). Politics "
        "does NOT dramatically underperform sports or crypto on sequential retention — if anything "
        f"crypto leads slightly at M12 ({_pct(m.m12_crypto)}). Category-first acquisition strategy "
        "won't move the needle much; habit infrastructure will.",
    )


def slide_election(prs: Presentation, m: DeckMetrics) -> None:
    slide = _blank_slide(prs)
    _headline(slide, "What drives stickiness is event MAGNITUDE, not category")
    _subtext(
        slide,
        f"The 2024 election cohort — supposed 'tourists' — became the stickiest: politics "
        f"{_pct(m.election_m12_politics)} at M12 vs sports {_pct(m.election_m12_sports)} / "
        f"crypto {_pct(m.election_m12_crypto)} on mature averages. Mega-events acquire high-intent users.",
    )
    _add_image(
        slide,
        _FIGURES / "03_election_cohort_retention.png",
        Inches(0.75),
        Inches(2.35),
        Inches(11.85),
        Inches(4.55),
    )
    _speaker_notes(
        slide,
        "Hypothesis flip: I assumed election-first users were disposable tourists. The Oct 2024 "
        f"cohort ({m.spike_election:,} new takers) tells the opposite story — politics-first users "
        f"hit {_pct(m.election_m12_politics)} at M12, nearly double the mature-cohort politics "
        f"average ({_pct(m.m12_politics)}). Mega-events don't just inflate top-of-funnel; they "
        "surface high-intent participants. The play is to retain spike cohorts through the NEXT "
        "mega-event, not to avoid politics acquisition.",
    )


def slide_mix(prs: Presentation, m: DeckMetrics) -> None:
    slide = _blank_slide(prs)
    _headline(
        slide,
        f"Politics drove {m.politics_pct * 100:.0f}% of all acquisition — "
        "spikes are the prize, not the problem",
    )
    politics_pct_display = f"{m.politics_pct * 100:.1f}%"
    _subtext(
        slide,
        f"{m.politics_users:,} of {m.total_users:,} all-time users entered via politics "
        f"({politics_pct_display}); capture the spikes, don't avoid them.",
    )
    _add_image(
        slide,
        _FIGURES / "05_acquisition_mix.png",
        Inches(0.75),
        Inches(2.35),
        Inches(11.85),
        Inches(4.55),
    )
    _speaker_notes(
        slide,
        f"First-trade category mix across all cohorts: politics {politics_pct_display} "
        f"({m.politics_users:,}), sports {m.sports_users:,}, crypto {m.crypto_users:,}, "
        f"other {m.other_users:,}. Politics is the largest entry point because mega-events "
        "are politics-heavy — that's feature, not bug. Growth should run retention sprints "
        "during spikes, not starve acquisition quality debates by category.",
    )


def slide_recommendation(prs: Presentation, m: DeckMetrics) -> None:
    slide = _blank_slide(prs)
    _headline(slide, "Build the retention layer that turns a second visit into a habit")

    levers = [
        (
            "Fast cross-category discovery on day 1 — crypto's daily cadence retains best "
            f"long-run (M12 {_pct(m.m12_crypto)}); recurring hooks build habit."
        ),
        (
            "Lifecycle re-engagement of dormant spike-cohorts timed to the NEXT mega-event "
            f"(e.g. reactivate Oct 2024 cohort before the next {_k(m.spike_recent)}-trader spike)."
        ),
        "Run each mega-event as a retention sprint, not just an acquisition campaign.",
    ]
    _add_bullets(slide, Inches(0.75), Inches(1.85), Inches(7.2), Inches(3.2), levers, font_size=17)

    # North Star callout box
    box = slide.shapes.add_shape(
        1,  # rectangle
        Inches(8.35),
        Inches(2.0),
        Inches(4.25),
        Inches(2.35),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_BOX_BG
    box.line.color.rgb = COLOR_BOX_BORDER
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p0 = tf.paragraphs[0]
    p0.text = "North Star"
    p0.font.name = FONT
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_POLITICS

    p1 = tf.add_paragraph()
    p1.text = f"Month-6 sequential retention"
    p1.font.name = FONT
    p1.font.size = Pt(13)
    p1.font.color.rgb = COLOR_TEXT

    p2 = tf.add_paragraph()
    p2.text = f"{_pct(m.north_star_m6)} today → target +5pp"
    p2.font.name = FONT
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT

    p3 = tf.add_paragraph()
    p3.text = (
        f"A 5pp lift on 200k/mo cohorts ≈ {m.lift_5pp_on_200k:,} retained traders/mo"
    )
    p3.font.name = FONT
    p3.font.size = Pt(12)
    p3.font.color.rgb = COLOR_MUTED

    _speaker_notes(
        slide,
        "Three levers, one metric. North Star is month-6 sequential retention because it "
        f"measures habit, not a one-time return ({_pct(m.ever_returned_pooled)} ever-returned "
        f"vs {_pct(m.seq_m6)} at M6). Today {_pct(m.north_star_m6)}; a 5 percentage-point lift "
        f"on 200k monthly cohorts = {m.lift_5pp_on_200k:,} incremental retained traders per month "
        "— material at Polymarket scale.",
    )


def slide_methodology(prs: Presentation, m: DeckMetrics) -> None:
    slide = _blank_slide(prs)
    _headline(slide, "Methodology & caveats")

    bullets = [
        "Data: on-chain Polygon trades via Dune curated tables (polymarket_polygon.market_trades). "
        "Taker = demand-side trader — NOT unique humans; bots/dust inflate counts (~5% of wallets "
        "= 75% of volume per Bloomberg).",
        "Cohort = calendar month of first trade. Two retention definitions: sequential month-N "
        f"(active in exactly month N; M6 median {_pct(m.seq_m6)}) vs ever-returned (≥2 active "
        f"months; pooled mature {_pct(m.ever_returned_pooled)}).",
        "Category = Polymarket market tags (~99.99% coverage on condition_id join) + keyword "
        "fallback on question/event name when tags empty.",
        "Volume deliberately avoided: ~25% of historical volume estimated wash trading (Columbia; "
        "peaked ~60% Dec 2024). Active traders is the honest North Star.",
    ]
    _add_bullets(slide, Inches(0.75), Inches(1.85), Inches(11.8), Inches(4.8), bullets, font_size=15)
    _speaker_notes(
        slide,
        "Be upfront in Q&A. Wallet ≠ human; directionally correct for cohort comparisons. "
        "Sequential vs ever-returned explains the 75% debate. Category tagging is hybrid tags+regex. "
        "We skipped volume because wash trading contaminates it — trader counts are cleaner for "
        "retention work.",
    )


def build_deck(metrics: DeckMetrics) -> Presentation:
    required_figures = [
        "01_acquisition_by_month.png",
        "02_retention_by_category.png",
        "03_election_cohort_retention.png",
        "04_the_75pct_myth.png",
        "05_acquisition_mix.png",
    ]
    for name in required_figures:
        path = _FIGURES / name
        if not path.is_file():
            raise FileNotFoundError(
                f"Missing {path}. Run: python scripts/03_build_charts.py"
            )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, metrics)
    slide_bluf(prs, metrics)
    slide_acquisition(prs, metrics)
    slide_75pct(prs, metrics)
    slide_category(prs, metrics)
    slide_election(prs, metrics)
    slide_mix(prs, metrics)
    slide_recommendation(prs, metrics)
    slide_methodology(prs, metrics)

    return prs


def main() -> int:
    if hasattr(sys.stdout, "reconfigure"):
        try:
            sys.stdout.reconfigure(encoding="utf-8")
        except (LookupError, OSError):
            pass

    metrics = load_metrics()
    prs = build_deck(metrics)

    _OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(_OUTPUT))

    print("=" * 60)
    print("Deck built from processed CSVs + figure PNGs")
    print("=" * 60)
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Saved:  {_OUTPUT}")
    print(f"  Author: {DECK_AUTHOR}  (set DECK_AUTHOR env to customize)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
